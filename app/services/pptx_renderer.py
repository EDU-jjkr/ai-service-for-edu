"""
PPTX Renderer Service
Renders LessonDeck JSON to PowerPoint (.pptx) files using python-pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from io import BytesIO
from typing import Optional
import logging

from app.models.lesson_schema import LessonDeck, Slide, SlideType
from app.services.template_manager import get_template_path
from app.services.stock_photo_service import StockPhotoService
from app.services.image_processor import ImageProcessor

logger = logging.getLogger(__name__)


class PPTXRenderer:
    """
    Renders LessonDeck objects to PowerPoint files.
    
    Process:
    1. Load template based on theme
    2. Create slides from lesson.slides
    3. Inject content into placeholders
    4. Fetch and embed images
    5. Return .pptx as BytesIO
    """
    
    def __init__(self, theme: str = "default"):
        """
        Initialize renderer with a specific theme.
        
        Args:
            theme: Theme name (e.g., "default", "science_nature")
        """
        self.theme = theme
        self.template_path = get_template_path(theme)
        self.stock_photo_service = StockPhotoService()
        self.image_processor = ImageProcessor()
        
        logger.info(f"PPTXRenderer initialized with theme: {theme}")
    
    async def render_lesson_deck(self, lesson: LessonDeck) -> BytesIO:
        """
        Main rendering function: Convert LessonDeck to PPTX file.
        
        Args:
            lesson: LessonDeck object with meta, structure, and slides
            
        Returns:
            BytesIO containing the complete .pptx file
        """
        logger.info(f"Rendering lesson: {lesson.meta.topic} ({len(lesson.slides)} slides)")
        
        # Load template
        prs = Presentation(self.template_path)
        
        # Add title slide
        self._create_title_slide(prs, lesson)
        
        # Add content slides
        for slide_data in sorted(lesson.slides, key=lambda s: s.order):
            await self._create_content_slide(prs, slide_data, lesson.meta.subject)
        
        # Save to BytesIO
        output = BytesIO()
        prs.save(output)
        output.seek(0)
        
        logger.info(f"✓ PPTX rendered successfully")
        return output
    
    def _create_title_slide(self, prs: Presentation, lesson: LessonDeck):
        """
        Create title slide with lesson metadata.
        
        Args:
            prs: Presentation object
            lesson: LessonDeck with metadata
        """
        # Use Title Slide layout (typically index 0)
        title_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_layout)
        
        # Set title
        title = slide.shapes.title
        if title:
            title.text = lesson.meta.topic
        
        # Set subtitle (if available)
        if len(slide.placeholders) > 1:
            subtitle = slide.placeholders[1]
            subtitle.text = f"Grade {lesson.meta.grade} - {lesson.meta.standards[0] if lesson.meta.standards else ''}"
        
        logger.info(f"✓ Title slide created")
    
    async def _create_content_slide(self, prs: Presentation, slide_data: Slide, subject: str):
        """
        Create a content slide with text and optional image.
        
        Args:
            prs: Presentation object
            slide_data: Slide object with content
            subject: Subject name for context
        """
        # Select layout based on slide type
        layout_idx = self._get_layout_index(slide_data.slideType)
        layout = prs.slide_layouts[layout_idx]
        slide = prs.slides.add_slide(layout)
        
        # Inject title
        if slide.shapes.title:
            slide.shapes.title.text = slide_data.title
        
        # Inject content
        self._inject_content(slide, slide_data.content)
        
        # Inject speaker notes
        if slide_data.speakerNotes:
            self._add_speaker_notes(slide, slide_data.speakerNotes)
        
        # Fetch and insert image if imageQuery exists
        if slide_data.imageQuery:
            await self._insert_image_from_query(slide, slide_data.imageQuery, slide_data.title)
        
        logger.info(f"✓ Slide created: {slide_data.title}")
    
    def _get_layout_index(self, slide_type: SlideType) -> int:
        """
        Map SlideType to layout index in template.
        
        Args:
            slide_type: SlideType enum
            
        Returns:
            Layout index (0-based)
        """
        # Default mapping for basic templates
        # In production, this would be configurable per theme
        mapping = {
            SlideType.INTRODUCTION: 0,  # Title slide layout
            SlideType.CONCEPT: 1,       # Title and Content
            SlideType.ACTIVITY: 1,      # Title and Content
            SlideType.ASSESSMENT: 1,    # Title and Content
            SlideType.SUMMARY: 1        # Title and Content
        }
        
        return mapping.get(slide_type, 1)
    
    def _inject_content(self, slide, content: str):
        """
        Inject text content into slide body placeholder.
        
        Args:
            slide: Slide object
            content: Text content (may be bulleted list)
        """
        # Find the body/content placeholder
        content_placeholder = None
        
        for shape in slide.placeholders:
            # Placeholder index 1 is typically the body/content
            if shape.placeholder_format.idx == 1:
                content_placeholder = shape
                break
        
        if not content_placeholder:
            logger.warning("No content placeholder found in slide")
            return
        
        # Set text
        text_frame = content_placeholder.text_frame
        text_frame.clear()  # Clear existing content
        
        # Split content into bullet points
        lines = content.strip().split('\n')
        
        for i, line in enumerate(lines):
            if not line.strip():
                continue
            
            if i == 0:
                # First paragraph already exists
                p = text_frame.paragraphs[0]
            else:
                # Add new paragraphs for subsequent lines
                p = text_frame.add_paragraph()
            
            p.text = line.strip()
            p.level = 0  # Bullet level
            p.font.size = Pt(18)
    
    def _add_speaker_notes(self, slide, notes: str):
        """
        Add speaker notes to slide.
        
        Args:
            slide: Slide object
            notes: Speaker notes text
        """
        if not slide.has_notes_slide:
            slide.notes_slide.notes_text_frame.text = notes
        else:
            slide.notes_slide.notes_text_frame.text += f"\n\n{notes}"
    
    async def _insert_image_from_query(self, slide, image_query: str, slide_title: str):
        """
        Fetch image from stock photo API and insert into slide.
        
        Args:
            slide: Slide object
            image_query: Search query for stock photos
            slide_title: Title of slide (for logging)
        """
        try:
            logger.info(f"Fetching image for: {image_query}")
            
            # Fetch image from stock photo service
            image_stream, attribution = await self.stock_photo_service.fetch_image(
                query=image_query,
                orientation="landscape"
            )
            
            if not image_stream:
                logger.warning(f"No image found for query: {image_query}")
                return
            
            # Find picture placeholder
            picture_placeholder = self._find_picture_placeholder(slide)
            
            if not picture_placeholder:
                logger.warning("No picture placeholder found, adding image to slide manually")
                # Add image at default position if no placeholder
                left = Inches(6)
                top = Inches(2)
                width = Inches(3)
                slide.shapes.add_picture(image_stream, left, top, width=width)
            else:
                # Crop image to fit placeholder aspect ratio
                ph_width = picture_placeholder.width
                ph_height = picture_placeholder.height
                
                # Convert EMU to pixels (approximate)
                target_width = int(ph_width / 9525)  # EMU to pixels conversion
                target_height = int(ph_height / 9525)
                
                cropped = self.image_processor.center_crop(
                    image_stream,
                    target_width=target_width,
                    target_height=target_height
                )
                
                # Compress for smaller file size
                compressed = self.image_processor.compress_for_pptx(
                    cropped,
                    max_size_kb=500
                )
                
                # Insert into placeholder
                picture_placeholder.insert_picture(compressed)
           
            # Add attribution to speaker notes
            if attribution:
                self._add_speaker_notes(slide, f"Image: {attribution}")
            
            logger.info(f"✓ Image inserted for: {slide_title}")
            
        except Exception as e:
            logger.error(f"Failed to insert image for '{slide_title}': {e}")
            # Continue without image if insertion fails
    
    def _find_picture_placeholder(self, slide):
        """
        Find the picture placeholder in a slide.
        
        Args:
            slide: Slide object
            
        Returns:
            Picture placeholder shape or None
        """
        for shape in slide.placeholders:
            # Type 18 is PP_PLACEHOLDER_TYPE.PICTURE
            if shape.placeholder_format.type == 18:
                return shape
        
        return None


# Convenience function
async def render_deck_to_pptx(lesson: LessonDeck, theme: str = "default") -> BytesIO:
    """
    Render a LessonDeck to PPTX file.
    
    Args:
        lesson: LessonDeck object
        theme: Theme name
        
    Returns:
        BytesIO containing .pptx file
    """
    renderer = PPTXRenderer(theme=theme)
    return await renderer.render_lesson_deck(lesson)
