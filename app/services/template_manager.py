"""
Template Manager
Manages PowerPoint template files and theme selection
"""

import os
from pathlib import Path
import logging

logger = logging.getLogger(__name__)


class TemplateManager:
    """
    Manages PowerPoint templates for different themes.
    
    Templates are stored in app/templates/ directory.
    Each template is a .pptx file with pre-designed master slides.
    """
    
    def __init__(self):
        # Get templates directory
        self.templates_dir = Path(__file__).parent.parent / "templates"
        
        # Ensure templates directory exists
        self.templates_dir.mkdir(exist_ok=True)
        
        logger.info(f"Template directory: {self.templates_dir}")
    
    def get_template_path(self, theme: str = "default") -> str:
        """
        Get the file path for a specific theme template.
        
        Args:
            theme: Theme name (e.g., "default", "science_nature", "mathematics")
            
        Returns:
            Absolute path to the template file
            
        Raises:
            FileNotFoundError: If template doesn't exist
        """
        template_file = self.templates_dir / f"{theme}.pptx"
        
        if not template_file.exists():
            logger.warning(f"Template '{theme}' not found, falling back to 'default'")
            
            # Fallback to default
            default_template = self.templates_dir / "default.pptx"
            
            if not default_template.exists():
                raise FileNotFoundError(
                    f"No templates found. Please create {default_template}"
                )
            
            return str(default_template.absolute())
        
        return str(template_file.absolute())
    
    def list_available_themes(self) -> list:
        """
        List all available template themes.
        
        Returns:
            List of theme names (without .pptx extension)
        """
        templates = list(self.templates_dir.glob("*.pptx"))
        themes = [t.stem for t in templates]
        
        logger.info(f"Available themes: {themes}")
        
        return themes
    
    def validate_template(self, template_path: str) -> bool:
        """
        Validate that a template file is a valid PowerPoint file.
        
        Args:
            template_path: Path to template file
            
        Returns:
            True if valid, False otherwise
        """
        try:
            from pptx import Presentation
            
            # Try to load the template
            prs = Presentation(template_path)
            
            # Check that it has at least one slide layout
            if len(prs.slide_layouts) == 0:
                logger.error(f"Template has no slide layouts: {template_path}")
                return False
            
            logger.info(f"Template valid: {len(prs.slide_layouts)} layouts found")
            return True
            
        except Exception as e:
            logger.error(f"Template validation failed: {e}")
            return False
    
    def get_layout_info(self, theme: str = "default") -> dict:
        """
        Get information about slide layouts in a template.
        
        Args:
            theme: Theme name
            
        Returns:
            Dict with layout names and placeholder counts
        """
        try:
            from pptx import Presentation
            
            template_path = self.get_template_path(theme)
            prs = Presentation(template_path)
            
            layout_info = {}
            
            for i, layout in enumerate(prs.slide_layouts):
                layout_info[i] = {
                    "name": layout.name,
                    "placeholders": len(layout.placeholders),
                    "placeholder_types": [
                        {
                            "idx": ph.placeholder_format.idx,
                            "type": ph.placeholder_format.type
                        }
                        for ph in layout.placeholders
                    ]
                }
            
            return layout_info
            
        except Exception as e:
            logger.error(f"Failed to get layout info: {e}")
            return {}


# Singleton instance
_template_manager = TemplateManager()


def get_template_path(theme: str = "default") -> str:
    """
    Convenience function to get template path.
    
    Args:
        theme: Theme name
        
    Returns:
        Absolute path to template file
    """
    return _template_manager.get_template_path(theme)


def list_available_themes() -> list:
    """
    Convenience function to list themes.
    
    Returns:
        List of available theme names
    """
    return _template_manager.list_available_themes()
