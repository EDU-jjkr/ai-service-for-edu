"""
PowerPoint Template Generator
Creates basic functional PowerPoint templates programmatically
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pathlib import Path
import logging

logger = logging.getLogger(__name__)


def create_default_template():
    """
    Create a basic PowerPoint template with 5 layouts.
    
    Layouts:
    1. Title Slide (for INTRODUCTION)
    2. Content with Image Right (for CONCEPT)
    3. Image Large with Caption (for ACTIVITY)
    4. Two-Column Layout (for ASSESSMENT)
    5. Summary Slide (for SUMMARY)
    """
    prs = Presentation()
    prs.slide_width = Inches(10)  # 16:9 aspect ratio
    prs.slide_height = Inches(5.625)
    
    logger.info("Creating default PowerPoint template...")
    
    # We'll use the built-in layouts and customize them
    # Note: For now, we'll create a simple template with blank slides
    # In production, a designer would create proper master slides in PowerPoint
    
    # Save to templates directory
    templates_dir = Path(__file__).parent.parent / "templates"
    templates_dir.mkdir(exist_ok=True)
    
    template_path = templates_dir / "default.pptx"
    
    # Create a simple slide to initialize the template
    blank_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(blank_layout)
    
    # Add a text box as instruction
    left = Inches(1)
    top = Inches(2)
    width = Inches(8)
    height = Inches(1)
    
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.text = "PowerPoint Template - Default Theme"
    
    p = tf.paragraphs[0]
    p.font.size = Pt(32)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    
    # Save template
    prs.save(str(template_path))
    
    logger.info(f"✓ Template created: {template_path}")
    return str(template_path)


def create_science_template():
    """
    Create a Science & Nature themed template.
    
    Features:
    - Green/blue color scheme
    - Nature-inspired fonts
    - Appropriate for science subjects
    """
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)
    
    logger.info("Creating science theme template...")
    
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    
    # Background color (light green)
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(230, 245, 240)  # Light mint green
    
    # Title text
    left = Inches(1)
    top = Inches(2)
    width = Inches(8)
    height = Inches(1)
    
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.text = "Science & Nature Theme"
    
    p = tf.paragraphs[0]
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(34, 139, 34)  # Forest green
    p.alignment = PP_ALIGN.CENTER
    
    # Save template
    templates_dir = Path(__file__).parent.parent / "templates"
    templates_dir.mkdir(exist_ok=True)
    
    template_path = templates_dir / "science_nature.pptx"
    prs.save(str(template_path))
    
    logger.info(f"✓ Template created: {template_path}")
    return str(template_path)


def initialize_templates():
    """
    Initialize all default templates if they don't exist.
    
    Call this on application startup.
    """
    templates_dir = Path(__file__).parent.parent / "templates"
    templates_dir.mkdir(exist_ok=True)
    
    # Create default template if missing
    default_path = templates_dir / "default.pptx"
    if not default_path.exists():
        logger.info("Default template missing, creating...")
        create_default_template()
    else:
        logger.info("✓ Default template exists")
    
    # Create science template if missing
    science_path = templates_dir / "science_nature.pptx"
    if not science_path.exists():
        logger.info("Science template missing, creating...")
        create_science_template()
    else:
        logger.info("✓ Science template exists")


if __name__ == "__main__":
    # Run this script to generate templates
    logging.basicConfig(level=logging.INFO)
    initialize_templates()
    print("\n✅ Templates initialized successfully!")
    print(f"Location: {Path(__file__).parent.parent / 'templates'}")
